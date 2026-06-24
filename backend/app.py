import os
import io
import json
import re
import tempfile
import time
import uuid
import zipfile
from datetime import datetime, timedelta, timezone
from urllib.parse import quote, unquote, urlparse
from urllib.request import Request, urlopen
from urllib.error import HTTPError, URLError
import numpy as np
from flask import Flask, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from werkzeug.exceptions import HTTPException
from PIL import Image
import tifffile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import base64

app = Flask(__name__)

FRONTEND_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "frontend"))
FRONTEND_FILES = {"index.html", "styles.css", "config.js", "app.js"}
FRONTEND_CSP = (
    "default-src 'self'; "
    "script-src 'self' https://cdn.sheetjs.com https://esm.sh; "
    "style-src 'self' 'unsafe-inline' https://fonts.googleapis.com; "
    "font-src 'self' https://fonts.gstatic.com; "
    "img-src 'self' blob: data:; "
    "connect-src 'self' https://blob.vercel-storage.com https://*.blob.vercel-storage.com "
    "http://127.0.0.1:* http://localhost:*; "
    "object-src 'none'; base-uri 'self'; frame-ancestors 'none'; form-action 'self'"
)


@app.errorhandler(Exception)
def handle_unexpected_error(error):
    if isinstance(error, HTTPException):
        return jsonify({"error": error.description}), error.code
    app.logger.exception("Unhandled backend error")
    return jsonify({"error": "Unexpected backend error."}), 500


@app.after_request
def add_security_headers(response):
    response.headers.setdefault("Cache-Control", "no-store")
    response.headers.setdefault("Content-Security-Policy", "default-src 'none'; frame-ancestors 'none'")
    response.headers.setdefault("Permissions-Policy", "camera=(), microphone=(), geolocation=()")
    response.headers.setdefault("Referrer-Policy", "no-referrer")
    response.headers.setdefault("Strict-Transport-Security", "max-age=31536000; includeSubDomains")
    response.headers.setdefault("X-Content-Type-Options", "nosniff")
    response.headers.setdefault("X-Frame-Options", "DENY")
    return response


def frontend_file_response(filename="index.html"):
    filename = filename.strip("/")
    if filename not in FRONTEND_FILES:
        return jsonify({"error": "Not found."}), 404
    response = send_from_directory(FRONTEND_DIR, filename)
    response.headers["Content-Security-Policy"] = FRONTEND_CSP
    return response


@app.route("/", methods=["GET"])
def frontend_index():
    return frontend_file_response("index.html")


@app.route("/frontend/<path:filename>", methods=["GET"])
def frontend_prefixed_file(filename):
    return frontend_file_response(filename)


@app.route("/<path:filename>", methods=["GET"])
def frontend_asset_file(filename):
    return frontend_file_response(filename)

# ─── CORS ─────────────────────────────────────────────────────────────────────
LOCAL_ORIGINS = [
    "http://localhost:8080",
    "http://127.0.0.1:8080",
    "http://[::1]:8080",
    "http://[::]:8080",
    "http://localhost:5500",
    "http://127.0.0.1:5500",
    "http://[::1]:5500",
]
def configured_allowed_origins():
    origins = [
        origin.strip().rstrip("/")
        for origin in os.environ.get("ALLOWED_ORIGINS", ",".join(LOCAL_ORIGINS)).split(",")
        if origin.strip()
    ]
    for origin in origins:
        parsed = urlparse(origin)
        if origin == "*" or parsed.scheme not in ("http", "https") or not parsed.netloc:
            raise RuntimeError("ALLOWED_ORIGINS must contain exact http(s) origins; wildcards are not allowed.")
        if parsed.path not in ("", "/") or parsed.params or parsed.query or parsed.fragment:
            raise RuntimeError("ALLOWED_ORIGINS entries must not contain paths, query strings, or fragments.")
    return origins


ALLOWED_ORIGINS = configured_allowed_origins()
CORS(
    app,
    origins=ALLOWED_ORIGINS,
    methods=["GET", "POST", "OPTIONS"],
    allow_headers=["Content-Type", "X-Blot-Session"],
    supports_credentials=False,
    max_age=600,
)

# ─── Rate limiting ─────────────────────────────────────────────────────────────
limiter = Limiter(
    get_remote_address,
    app=app,
    default_limits=["200 per day", "50 per hour"],
    storage_uri=os.environ.get("RATELIMIT_STORAGE_URI", "memory://"),
)

# ─── Constants ────────────────────────────────────────────────────────────────
MAX_ZIP_BYTES = int(os.environ.get("MAX_ZIP_BYTES", 250 * 1024 * 1024))
MAX_TIF_BYTES = int(os.environ.get("MAX_TIF_BYTES", 100 * 1024 * 1024))
MAX_ZIP_ENTRIES = 400
MAX_ZIP_UNCOMPRESSED = int(os.environ.get("MAX_ZIP_UNCOMPRESSED_BYTES", 400 * 1024 * 1024))
MAX_ZIP_COMPRESSION_RATIO = 200
MAX_IMAGE_PIXELS = 80_000_000
MAX_TEXT_BYTES = 2 * 1024 * 1024
MAX_JPEG_BYTES = 50 * 1024 * 1024
MAX_NAME_LENGTH = 200
MAX_TIF_PAGES = 16
MAX_REQUEST_BYTES = int(os.environ.get("MAX_REQUEST_BYTES", 16 * 1024 * 1024))
MAX_PPTX_SLIDES = 40
MAX_PPTX_GRAPHS = 120
MAX_PPTX_IMAGE_BYTES = 8 * 1024 * 1024
MAX_PPTX_TOTAL_IMAGE_BYTES = 40 * 1024 * 1024
BLOT_TIMEZONE_OFFSETS = {
    "UTC": 0,
    "GMT": 0,
    "PST": -8 * 60,
    "PDT": -7 * 60,
    "MST": -7 * 60,
    "MDT": -6 * 60,
    "CST": -6 * 60,
    "CDT": -5 * 60,
    "EST": -5 * 60,
    "EDT": -4 * 60,
    "AKST": -9 * 60,
    "AKDT": -8 * 60,
    "HST": -10 * 60,
    "AST": -4 * 60,
    "ADT": -3 * 60,
    "NST": -(3 * 60 + 30),
    "NDT": -(2 * 60 + 30),
    "CET": 1 * 60,
    "CEST": 2 * 60,
}
BLOT_TIMESTAMP_PATTERN = re.compile(
    r"^#(?P<weekday>Mon|Tue|Wed|Thu|Fri|Sat|Sun)\s+"
    r"(?P<month>Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\s+"
    r"(?P<day>\d{1,2})\s+(?P<clock>\d{2}:\d{2}:\d{2})\s+"
    r"(?P<zone>[A-Za-z]{2,5}|[+-]\d{4})\s+(?P<year>\d{4})$"
)
BLOT_FILE_FIELDS = {
    "jpg": ("has_jpg", "jpg_bytes", "preview.jpg"),
    "700": ("has_700", "tif_700_bytes", "700.tif"),
    "800": ("has_800", "tif_800_bytes", "800.tif"),
}
ALL_BLOT_FILE_KINDS = tuple(BLOT_FILE_FIELDS)
BLOT_FILE_LIMITS = {
    "jpg": MAX_JPEG_BYTES,
    "700": MAX_TIF_BYTES,
    "800": MAX_TIF_BYTES,
}
SAFE_PATH_PART_PATTERN = re.compile(r"^[A-Za-z0-9._-]{1,200}$")
SESSION_FILE_NAMES = {field[2] for field in BLOT_FILE_FIELDS.values()}
TEMP_STORAGE_BACKEND = os.environ.get("BLOT_TEMP_STORAGE", "local").lower().replace("_", "-")
LOCAL_TEMP_DIR = os.environ.get("BLOT_TEMP_DIR") or os.path.join(
    tempfile.gettempdir(),
    "western-blot-browser-only",
)
USE_VERCEL_BLOB = TEMP_STORAGE_BACKEND == "vercel-blob"
BLOB_ACCESS = "public" if os.environ.get("BLOB_ACCESS") == "public" else "private"
BLOB_API_BASE_URL = "https://blob.vercel-storage.com"
BLOB_API_VERSION = "10"
app.config["MAX_CONTENT_LENGTH"] = MAX_REQUEST_BYTES

def init_storage():
    if USE_VERCEL_BLOB:
        if not os.environ.get("BLOB_READ_WRITE_TOKEN"):
            raise RuntimeError("BLOB_READ_WRITE_TOKEN is required when BLOT_TEMP_STORAGE=vercel-blob.")
        return
    os.makedirs(LOCAL_TEMP_DIR, exist_ok=True)


def now_iso():
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def parse_blot_created_at(value):
    match = BLOT_TIMESTAMP_PATTERN.fullmatch(str(value or "").strip())
    if not match:
        return None

    zone_name = match.group("zone").upper()
    if re.fullmatch(r"[+-]\d{4}", zone_name):
        sign = 1 if zone_name[0] == "+" else -1
        offset_minutes = sign * (int(zone_name[1:3]) * 60 + int(zone_name[3:5]))
    else:
        offset_minutes = BLOT_TIMEZONE_OFFSETS.get(zone_name)
    if offset_minutes is None or abs(offset_minutes) > 14 * 60:
        return None

    try:
        parsed = datetime.strptime(
            "{weekday} {month} {day} {clock} {year}".format(**match.groupdict()),
            "%a %b %d %H:%M:%S %Y",
        )
    except ValueError:
        return None

    captured_at = parsed.replace(tzinfo=timezone(timedelta(minutes=offset_minutes)))
    return captured_at.astimezone(timezone.utc).isoformat(timespec="seconds")


def safe_id(value):
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", str(value))
    return cleaned.strip("._") or uuid.uuid4().hex


def request_session_id(data=None):
    if data is None:
        data = {}
    candidate = (
        request.headers.get("X-Blot-Session")
        or request.form.get("sessionId")
        or data.get("sessionId")
        or data.get("session_id")
        or ""
    )
    return safe_id(candidate)[:80] if candidate else uuid.uuid4().hex


def temp_object_path(session_id, blot_id, filename):
    return f"sessions/{safe_id(session_id)}/{safe_id(blot_id)}/{safe_id(filename)}"


def validate_temp_path(path, session_id=None, allow_uploads=False):
    path = str(path or "")
    if len(path) > 500 or not path:
        raise PublicError("Invalid temporary file reference.")
    if "\\" in path or path.startswith("/") or "\x00" in path:
        raise PublicError("Invalid temporary file reference.")
    parts = path.split("/")
    if any(part in ("", ".", "..") for part in parts):
        raise PublicError("Invalid temporary file reference.")
    if any(not SAFE_PATH_PART_PATTERN.fullmatch(part) for part in parts):
        raise PublicError("Invalid temporary file reference.")
    allowed_roots = {"sessions"}
    if allow_uploads:
        allowed_roots.add("uploads")
    if parts[0] not in allowed_roots:
        raise PublicError("Invalid temporary file reference.")
    if parts[0] == "sessions":
        if len(parts) != 4 or parts[3] not in SESSION_FILE_NAMES:
            raise PublicError("Invalid temporary file reference.")
    if parts[0] == "uploads":
        if len(parts) != 3 or not parts[2].lower().endswith(".zip"):
            raise PublicError("Invalid temporary file reference.")
    if session_id and len(parts) > 1 and parts[1] != safe_id(session_id):
        raise PublicError("Temporary file does not belong to this browser session.", 403)
    return "/".join(parts)


def descriptor_path(descriptor, session_id=None, allow_uploads=False):
    if not isinstance(descriptor, dict):
        raise PublicError("Invalid temporary file reference.")
    path = descriptor.get("pathname") or descriptor.get("path")
    return validate_temp_path(path, session_id, allow_uploads=allow_uploads)


def local_temp_file_path(path):
    path = validate_temp_path(path, allow_uploads=True)
    full_path = os.path.abspath(os.path.join(LOCAL_TEMP_DIR, *path.split("/")))
    root = os.path.abspath(LOCAL_TEMP_DIR)
    if not full_path.startswith(root + os.sep):
        raise PublicError("Invalid temporary file reference.")
    return full_path


def descriptor_from_blob_result(result, fallback_path, content_type):
    result = result or {}
    return {
        "backend": "vercel-blob",
        "path": result.get("pathname") or result.get("path") or fallback_path,
        "pathname": result.get("pathname") or result.get("path") or fallback_path,
        "url": result.get("url"),
        "downloadUrl": result.get("downloadUrl") or result.get("download_url"),
        "contentType": result.get("contentType") or result.get("content_type") or content_type,
    }


def vercel_blob_token():
    token = os.environ.get("BLOB_READ_WRITE_TOKEN", "")
    if not token:
        raise RuntimeError("BLOB_READ_WRITE_TOKEN is required when BLOT_TEMP_STORAGE=vercel-blob.")
    return token


def is_allowed_vercel_blob_host(hostname):
    host = (hostname or "").lower().rstrip(".")
    return host == "blob.vercel-storage.com" or host.endswith(".blob.vercel-storage.com")


def validate_vercel_blob_url(url, expected_path=None):
    if not isinstance(url, str):
        raise PublicError("Invalid temporary file reference.")
    parsed = urlparse(url)
    if parsed.scheme != "https" or not is_allowed_vercel_blob_host(parsed.hostname):
        raise PublicError("Invalid temporary file reference.")
    if expected_path is not None:
        decoded_path = unquote(parsed.path.lstrip("/"))
        if decoded_path != expected_path:
            raise PublicError("Temporary file URL does not match its storage path.", 403)
    return url


def vercel_blob_request(method, url, body=None, headers=None, timeout=60, max_response_bytes=None):
    validate_vercel_blob_url(url)
    request_headers = {
        "authorization": f"Bearer {vercel_blob_token()}",
        "x-api-version": BLOB_API_VERSION,
    }
    if headers:
        request_headers.update(headers)
    request = Request(url, data=body, headers=request_headers, method=method)
    try:
        with urlopen(request, timeout=timeout) as response:
            if max_response_bytes is None:
                response_body = response.read()
            else:
                response_body = response.read(max_response_bytes + 1)
                if len(response_body) > max_response_bytes:
                    raise PublicError("Temporary file exceeds the configured size limit.", 413)
            content_type = response.headers.get("Content-Type", "")
    except (HTTPError, URLError, TimeoutError) as error:
        raise PublicError("Temporary Blob storage request failed.", 502) from error
    if not response_body:
        return None
    if "application/json" in content_type:
        return json.loads(response_body.decode("utf-8"))
    return response_body


def vercel_blob_put(path, file_bytes, content_type):
    encoded_path = quote(path, safe="/")
    return vercel_blob_request(
        "PUT",
        f"{BLOB_API_BASE_URL}/?pathname={encoded_path}",
        body=file_bytes,
        headers={
            "access": BLOB_ACCESS,
            "x-content-type": content_type,
            "x-cache-control-max-age": "60",
            "x-allow-overwrite": "0",
        },
    )


def vercel_blob_read(descriptor, path, max_bytes=None):
    url = descriptor.get("downloadUrl") or descriptor.get("url")
    validate_vercel_blob_url(url, path)
    return vercel_blob_request("GET", url, max_response_bytes=max_bytes)


def vercel_blob_delete(descriptor_paths):
    urls = []
    for descriptor, path in descriptor_paths:
        if not isinstance(descriptor, dict) or not isinstance(descriptor.get("url"), str):
            continue
        try:
            urls.append(validate_vercel_blob_url(descriptor["url"], path))
        except PublicError:
            continue
    if not urls:
        return
    body = json.dumps({"urls": urls}).encode("utf-8")
    vercel_blob_request(
        "POST",
        f"{BLOB_API_BASE_URL}/delete",
        body=body,
        headers={"Content-Type": "application/json"},
    )


def store_temp_file(session_id, blot_id, filename, file_bytes, content_type):
    if not file_bytes:
        return None
    path = temp_object_path(session_id, blot_id, filename)
    if USE_VERCEL_BLOB:
        result = vercel_blob_put(path, file_bytes, content_type)
        return descriptor_from_blob_result(result, path, content_type)

    full_path = local_temp_file_path(path)
    os.makedirs(os.path.dirname(full_path), exist_ok=True)
    with open(full_path, "wb") as handle:
        handle.write(file_bytes)
    return {
        "backend": "local",
        "path": path,
        "pathname": path,
        "contentType": content_type,
    }


def read_temp_file(descriptor, session_id=None, allow_uploads=False, max_bytes=None):
    path = descriptor_path(descriptor, session_id, allow_uploads=allow_uploads)
    if USE_VERCEL_BLOB:
        return vercel_blob_read(descriptor, path, max_bytes=max_bytes)

    full_path = local_temp_file_path(path)
    if not os.path.exists(full_path):
        raise PublicError("Temporary file was not found.", 404)
    with open(full_path, "rb") as handle:
        if max_bytes is None:
            return handle.read()
        file_bytes = handle.read(max_bytes + 1)
    if len(file_bytes) > max_bytes:
        raise PublicError("Temporary file exceeds the configured size limit.", 413)
    return file_bytes


def delete_temp_files(descriptors, session_id=None, allow_uploads=False):
    checked_descriptors = []
    for descriptor in descriptors:
        try:
            path = descriptor_path(descriptor, session_id, allow_uploads=allow_uploads)
            checked_descriptors.append((descriptor, path))
        except PublicError:
            continue
    if not checked_descriptors:
        return
    if USE_VERCEL_BLOB:
        vercel_blob_delete(checked_descriptors)
        return
    for _descriptor, path in checked_descriptors:
        full_path = local_temp_file_path(path)
        if os.path.exists(full_path):
            os.remove(full_path)


def collect_blot_file_descriptors(blots):
    descriptors = []
    for blot in blots or []:
        files = blot.get("files", {}) if isinstance(blot, dict) else {}
        for descriptor in files.values():
            if descriptor:
                descriptors.append(descriptor)
    return descriptors


def load_payload_blot(blot, file_kinds=ALL_BLOT_FILE_KINDS, session_id=None):
    if not isinstance(blot, dict):
        raise PublicError("Blot data is missing.")
    files = blot.get("files")
    if not isinstance(files, dict):
        raise PublicError("Blot temporary file references are missing.")
    loaded = {"id": blot.get("id"), "name": blot.get("name", "Untitled blot"), "files": files}
    for kind in file_kinds:
        descriptor = files.get(kind)
        bytes_field = BLOT_FILE_FIELDS[kind][1]
        loaded[bytes_field] = (
            read_temp_file(descriptor, session_id, max_bytes=BLOT_FILE_LIMITS[kind])
            if descriptor
            else None
        )
    return loaded


init_storage()

# ─── Health check ─────────────────────────────────────────────────────────────

@app.route("/health")
@app.route("/api/health")
def health():
    return jsonify({"status": "ok"})


@app.route("/client-config")
@app.route("/api/client-config")
def client_config():
    return jsonify({"blobAccess": BLOB_ACCESS})

# ─── ZIP upload & parsing ─────────────────────────────────────────────────────

@app.route("/upload-zip", methods=["POST"])
@app.route("/api/upload-zip", methods=["POST"])
@limiter.limit("10 per minute")
def upload_zip():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files["file"]
    if not file.filename.lower().endswith(".zip"):
        return jsonify({"error": "File must be a .zip"}), 400

    # Check content length before reading
    if request.content_length and request.content_length > MAX_ZIP_BYTES:
        return jsonify({"error": "File is too large for direct API upload. Use Storage upload instead."}), 413

    file_bytes = file.read()

    # Check actual file size after reading
    if len(file_bytes) > MAX_ZIP_BYTES:
        return jsonify({"error": "File is too large for direct API upload. Use Storage upload instead."}), 413

    # Validate ZIP magic bytes
    if not is_valid_zip(file_bytes):
        return jsonify({"error": "Invalid ZIP file."}), 400

    try:
        blots = parse_zip(file_bytes, request_session_id())
        return jsonify({"blots": blots})
    except PublicError as error:
        return error_response(error, "ZIP import failed.")
    except Exception as error:
        app.logger.exception("ZIP upload failed")
        return error_response(error, "ZIP import failed.")


@app.route("/process-upload", methods=["POST"])
@app.route("/api/process-upload", methods=["POST"])
@limiter.limit("10 per minute")
def process_storage_upload():
    data = request.get_json(silent=True) or {}
    session_id = request_session_id(data)
    upload_descriptor = data.get("upload") or data.get("blob") or {}
    started = time.monotonic()
    upload_path = upload_descriptor.get("pathname") or upload_descriptor.get("path")
    print(
        json.dumps({
            "event": "process_upload_start",
            "sessionId": session_id,
            "path": upload_path,
        }),
        flush=True,
    )
    try:
        file_bytes = read_temp_file(upload_descriptor, session_id, allow_uploads=True, max_bytes=MAX_ZIP_BYTES)
        print(
            json.dumps({
                "event": "process_upload_blob_read",
                "sessionId": session_id,
                "bytes": len(file_bytes),
                "elapsedSeconds": round(time.monotonic() - started, 3),
            }),
            flush=True,
        )
        if not is_valid_zip(file_bytes):
            raise PublicError("Stored object is not a valid ZIP file.")
        blots = parse_zip(file_bytes, session_id)
        print(
            json.dumps({
                "event": "process_upload_zip_parsed",
                "sessionId": session_id,
                "blotCount": len(blots),
                "elapsedSeconds": round(time.monotonic() - started, 3),
            }),
            flush=True,
        )
        delete_temp_files([upload_descriptor], session_id, allow_uploads=True)
        print(
            json.dumps({
                "event": "process_upload_done",
                "sessionId": session_id,
                "elapsedSeconds": round(time.monotonic() - started, 3),
            }),
            flush=True,
        )
        return jsonify({"blots": blots})
    except PublicError as error:
        print(
            json.dumps({
                "event": "process_upload_public_error",
                "sessionId": session_id,
                "error": str(error),
                "elapsedSeconds": round(time.monotonic() - started, 3),
            }),
            flush=True,
        )
        return error_response(error, "ZIP import failed.")
    except Exception as error:
        print(
            json.dumps({
                "event": "process_upload_unexpected_error",
                "sessionId": session_id,
                "elapsedSeconds": round(time.monotonic() - started, 3),
            }),
            flush=True,
        )
        app.logger.exception("Stored ZIP processing failed")
        return error_response(error, "ZIP import failed.")


def is_valid_zip(file_bytes):
    return len(file_bytes) >= 4 and zipfile.is_zipfile(io.BytesIO(file_bytes))


def enforce_zip_member_size(zf, filename, limit, message):
    """Reject an optional archive member before reading oversized data into memory."""
    if filename and zf.getinfo(filename).file_size > limit:
        raise PublicError(message, 413)


def parse_zip(file_bytes, session_id):
    blots = []
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
        validate_zip_archive(zf)
        # Group files by folder
        folders = {}
        for name in zf.namelist():
            if name.endswith("/"):
                continue
            parts = name.split("/")
            folder = parts[0] if len(parts) > 1 else "__root__"
            folders.setdefault(folder, []).append(name)

        for folder_index, (folder, files) in enumerate(folders.items(), start=1):
            txt_file = next((f for f in files if f.lower().endswith(".txt")), None)
            jpg_file = next((f for f in files if f.lower().endswith(".jpg") or f.lower().endswith(".jpeg")), None)
            tif_700  = next((f for f in files if "700" in f and f.lower().endswith(".tif")), None)
            tif_800  = next((f for f in files if "800" in f and f.lower().endswith(".tif")), None)

            if not txt_file:
                continue

            print(
                json.dumps({
                    "event": "process_upload_parse_folder",
                    "sessionId": session_id,
                    "folderIndex": folder_index,
                    "has700": tif_700 is not None,
                    "has800": tif_800 is not None,
                    "hasJpg": jpg_file is not None,
                }),
                flush=True,
            )

            # Enforce per-file limits before loading archive members into memory.
            enforce_zip_member_size(zf, tif_700, MAX_TIF_BYTES, "A 700nm TIF exceeds the configured size limit.")
            enforce_zip_member_size(zf, tif_800, MAX_TIF_BYTES, "An 800nm TIF exceeds the configured size limit.")
            enforce_zip_member_size(zf, txt_file, MAX_TEXT_BYTES, "Blot metadata text is too large.")
            enforce_zip_member_size(zf, jpg_file, MAX_JPEG_BYTES, "Blot preview image is too large.")

            # Parse blot name from last line of txt file
            txt_content = zf.read(txt_file).decode("utf-8", errors="ignore")
            lines = txt_content.strip().splitlines()
            last_line = lines[-1] if lines else ""
            blot_name = last_line.split("=", 1)[1].strip() if last_line.startswith("Remarks=") else folder
            blot_name = blot_name[:MAX_NAME_LENGTH] or "Untitled blot"
            created_at = parse_blot_created_at(lines[1] if len(lines) > 1 else None) or now_iso()

            # Generate a unique unpredictable ID
            blot_id = f"{uuid.uuid4().hex}_{folder}".replace(" ", "_")

            jpg_bytes = zf.read(jpg_file) if jpg_file else None
            tif_700_bytes = zf.read(tif_700) if tif_700 else None
            tif_800_bytes = zf.read(tif_800) if tif_800 else None
            validate_tif_pixels(tif_700_bytes, "700nm TIF")
            validate_tif_pixels(tif_800_bytes, "800nm TIF")

            files = {
                "jpg": store_temp_file(session_id, blot_id, "preview.jpg", jpg_bytes, "image/jpeg"),
                "700": store_temp_file(session_id, blot_id, "700.tif", tif_700_bytes, "image/tiff"),
                "800": store_temp_file(session_id, blot_id, "800.tif", tif_800_bytes, "image/tiff"),
            }
            files = {key: value for key, value in files.items() if value}
            blots.append({
                "id": blot_id,
                "name": blot_name,
                "hasJpg": jpg_file is not None,
                "has700": tif_700 is not None,
                "has800": tif_800 is not None,
                "scanCount": 0,
                "createdAt": created_at,
                "files": files,
            })

    return blots


class PublicError(Exception):
    def __init__(self, message, status=400):
        super().__init__(message)
        self.status = status


def error_response(error, fallback):
    """Return approved client errors while hiding unexpected exception details."""
    if isinstance(error, PublicError):
        return jsonify({"error": str(error)}), error.status
    return jsonify({"error": fallback}), 500


def validate_zip_archive(zf):
    infos = [info for info in zf.infolist() if not info.is_dir()]
    if len(infos) > MAX_ZIP_ENTRIES:
        raise PublicError(f"ZIP has too many files. Maximum is {MAX_ZIP_ENTRIES}.")

    total_size = sum(info.file_size for info in infos)
    if total_size > MAX_ZIP_UNCOMPRESSED:
        raise PublicError("ZIP contents are too large after decompression.")

    seen_paths = set()
    for info in infos:
        filename = info.filename
        if not filename or len(filename) > 500 or "\x00" in filename or "\\" in filename:
            raise PublicError("ZIP contains an invalid file name.")
        path_parts = filename.split("/")
        if os.path.isabs(filename) or any(part in ("", ".", "..") for part in path_parts):
            raise PublicError("ZIP contains unsafe file paths.")
        normalized = filename.casefold()
        if normalized in seen_paths:
            raise PublicError("ZIP contains duplicate file paths.")
        seen_paths.add(normalized)
        if info.flag_bits & 0x1:
            raise PublicError("Encrypted ZIP entries are not supported.")
        if info.file_size > 0:
            ratio = info.file_size / max(1, info.compress_size)
            if ratio > MAX_ZIP_COMPRESSION_RATIO:
                raise PublicError("ZIP contains a suspiciously compressed file.")


def validate_tif_pixels(tif_bytes, label):
    if not tif_bytes:
        return
    decode_validated_tif(tif_bytes, label)


def decode_validated_tif(tif_bytes, label):
    try:
        with tifffile.TiffFile(io.BytesIO(tif_bytes)) as tif:
            pages = list(tif.pages)
            if not pages or len(pages) > MAX_TIF_PAGES:
                raise PublicError(f"{label} has an unsupported number of image pages.")

            full_resolution_pages = [page for page in pages if not page.is_reduced]
            if len(full_resolution_pages) != 1:
                raise PublicError(
                    f"{label} must contain one full-resolution image; "
                    "additional pages may only be reduced-resolution previews."
                )

            primary_page = full_resolution_pages[0]
            primary_shape = primary_page.shape
            if len(primary_shape) != 2:
                raise PublicError(f"{label} must be a two-dimensional grayscale image.")

            supported_dtypes = {np.dtype("uint16"), np.dtype("float16")}
            total_pixels = 0
            for page in pages:
                if len(page.shape) != 2:
                    raise PublicError(f"{label} contains a non-grayscale image page.")
                if page.dtype not in supported_dtypes:
                    raise PublicError(f"{label} must contain 16-bit grayscale image data.")
                page_pixels = int(np.prod(page.shape))
                if page_pixels <= 0:
                    raise PublicError(f"{label} has invalid dimensions.")
                total_pixels += page_pixels
                if page.is_reduced and (
                    page.shape[0] > primary_shape[0] or page.shape[1] > primary_shape[1]
                ):
                    raise PublicError(f"{label} contains an invalid reduced-resolution preview.")

            pixels = int(np.prod(primary_shape))
            if pixels > MAX_IMAGE_PIXELS:
                raise PublicError(f"{label} is too large. Maximum is {MAX_IMAGE_PIXELS:,} pixels.")
            if total_pixels > MAX_IMAGE_PIXELS * 2:
                raise PublicError(f"{label} contains too much decoded image data.")

            image = primary_page.asarray()
            if image.shape != primary_shape:
                raise PublicError(f"{label} contains invalid pixel data.")
            return sanitize_tif_pixels(image, label)
    except PublicError:
        raise
    except Exception as error:
        raise PublicError(f"{label} could not be read.") from error


def sanitize_tif_pixels(image, label):
    if image.dtype != np.dtype("float16"):
        return image

    finite = np.isfinite(image)
    if not finite.any():
        raise PublicError(f"{label} does not contain any finite pixel measurements.")
    if finite.all():
        return image

    nan_count = int(np.isnan(image).sum())
    posinf_count = int(np.isposinf(image).sum())
    neginf_count = int(np.isneginf(image).sum())
    app.logger.info(
        "%s contains LI-COR special pixels: %d masked, %d saturated, %d negative-overflow",
        label,
        nan_count,
        posinf_count,
        neginf_count,
    )
    limit = np.finfo(np.float16).max
    return np.nan_to_num(image, copy=True, nan=0.0, posinf=limit, neginf=-limit)


# ─── TIF composite rendering ──────────────────────────────────────────────────

@app.route("/render-composite", methods=["POST"])
@app.route("/api/render-composite", methods=["POST"])
@limiter.limit("60 per minute")
def render_composite():
    data = request.get_json(silent=True) or {}
    session_id = request_session_id(data)
    try:
        blot = load_payload_blot(data.get("blot"), ("700", "800"), session_id)
        if not blot["tif_700_bytes"] or not blot["tif_800_bytes"]:
            return jsonify({"error": "TIF files not found for this blot"}), 404

        brightness_700 = float(data.get("brightness700", 1.0))
        contrast_700 = float(data.get("contrast700", 1.0))
        brightness_800 = float(data.get("brightness800", 1.0))
        contrast_800 = float(data.get("contrast800", 1.0))
        color_mode = data.get("colorMode", "color")

        brightness_700 = max(0.1, min(5.0, brightness_700))
        contrast_700 = max(0.1, min(10.0, contrast_700))
        brightness_800 = max(0.1, min(5.0, brightness_800))
        contrast_800 = max(0.1, min(10.0, contrast_800))

        img = build_composite(
            blot["tif_700_bytes"],
            blot["tif_800_bytes"],
            brightness_700, contrast_700,
            brightness_800, contrast_800,
            color_mode,
        )

        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=92)
        buf.seek(0)
        return send_file(buf, mimetype="image/jpeg")
    except PublicError as error:
        return error_response(error, "Could not render blot image.")
    except Exception as error:
        app.logger.exception("Composite render failed")
        return error_response(error, "Could not render blot image.")


def read_tif_channel(tif_bytes):
    # Display normalization starts from the same validated pixels used for quantification.
    arr = read_raw_channel(tif_bytes)
    low  = np.percentile(arr, 1)
    high = np.percentile(arr, 99)
    if high > low:
        arr = (arr - low) / (high - low)
    else:
        arr = np.zeros_like(arr)
    return np.clip(arr, 0.0, 1.0)


def apply_adjustments(channel, brightness, contrast):
    adjusted = (channel - 0.5) * contrast + 0.5
    adjusted = adjusted + (brightness - 1.0)
    return np.clip(adjusted, 0.0, 1.0)


def build_composite(tif_700_bytes, tif_800_bytes, brightness_700, contrast_700, brightness_800, contrast_800, color_mode="color"):
    ch_700 = read_tif_channel(tif_700_bytes)
    ch_800 = read_tif_channel(tif_800_bytes)

    if ch_700.shape != ch_800.shape:
        h = max(ch_700.shape[0], ch_800.shape[0])
        w = max(ch_700.shape[1], ch_800.shape[1])
        ch_700 = np.array(Image.fromarray((ch_700 * 65535).astype(np.uint16)).resize((w, h), Image.LANCZOS)).astype(np.float32) / 65535
        ch_800 = np.array(Image.fromarray((ch_800 * 65535).astype(np.uint16)).resize((w, h), Image.LANCZOS)).astype(np.float32) / 65535

    ch_700 = apply_adjustments(ch_700, brightness_700, contrast_700)
    ch_800 = apply_adjustments(ch_800, brightness_800, contrast_800)

    if color_mode == "grayscale":
        blended = np.clip(ch_700 + ch_800, 0.0, 1.0)
        blended = 1.0 - blended
        gray = (blended * 255).astype(np.uint8)
        rgb = np.stack([gray, gray, gray], axis=2)
    else:
        r = (ch_700 * 255).astype(np.uint8)
        g = (ch_800 * 255).astype(np.uint8)
        b = np.zeros_like(r)
        rgb = np.stack([r, g, b], axis=2)

    return Image.fromarray(rgb, mode="RGB")

# ─── Signal extraction ────────────────────────────────────────────────────────

@app.route("/extract", methods=["POST"])
@app.route("/api/extract", methods=["POST"])
@limiter.limit("30 per minute")
def extract_payload_signals():
    data = request.get_json(silent=True) or {}
    session_id = request_session_id(data)
    try:
        blot = load_payload_blot(data.get("blot"), (), session_id)
        boxes = data.get("boxes", [])
        channel = data.get("channel", "700")
        background_axis = data.get("backgroundAxis", "leftright")

        if channel not in ("700", "800"):
            return jsonify({"error": "Invalid channel"}), 400
        if background_axis not in ("leftright", "topbottom"):
            return jsonify({"error": "Invalid background mode"}), 400
        if not isinstance(boxes, list):
            return jsonify({"error": "Boxes must be a list."}), 400
        if len(boxes) > 200:
            return jsonify({"error": "Too many boxes. Maximum is 200."}), 400

        for box in boxes:
            validate_box(box)
        descriptor = blot["files"].get(channel)
        if not descriptor:
            return jsonify({"error": f"No {channel}nm TIF found"}), 404
        tif_bytes = read_temp_file(descriptor, session_id, max_bytes=BLOT_FILE_LIMITS[channel])
        arr = read_raw_channel(tif_bytes)
        results = [extract_box_signal(arr, box, background_axis) for box in boxes]
        return jsonify({"results": results})
    except PublicError as error:
        return error_response(error, "Signal extraction failed.")
    except Exception as error:
        app.logger.exception("Signal extraction failed")
        return error_response(error, "Signal extraction failed.")


def read_raw_channel(tif_bytes):
    return read_validated_tif(tif_bytes).astype(np.float32)


def read_validated_tif(tif_bytes):
    return decode_validated_tif(tif_bytes, "Stored TIF")


def extract_box_signal(arr, box, background_axis):
    x = max(0, int(round(box["x"])))
    y = max(0, int(round(box["y"])))
    w = max(1, int(round(box["w"])))
    h = max(1, int(round(box["h"])))

    img_h, img_w = arr.shape
    x2 = min(x + w, img_w)
    y2 = min(y + h, img_h)

    roi        = arr[y:y2, x:x2]
    raw_signal = float(np.sum(roi))

    border     = 3
    bg_pixels  = []

    if background_axis == "leftright":
        lx1 = max(0, x - border)
        lx2 = x
        if lx2 > lx1:
            bg_pixels.append(arr[y:y2, lx1:lx2].flatten())
        rx1 = x2
        rx2 = min(img_w, x2 + border)
        if rx2 > rx1:
            bg_pixels.append(arr[y:y2, rx1:rx2].flatten())
    else:
        ty1 = max(0, y - border)
        ty2 = y
        if ty2 > ty1:
            bg_pixels.append(arr[ty1:ty2, x:x2].flatten())
        by1 = y2
        by2 = min(img_h, y2 + border)
        if by2 > by1:
            bg_pixels.append(arr[by1:by2, x:x2].flatten())

    if bg_pixels:
        all_bg     = np.concatenate(bg_pixels)
        bg_median  = float(np.median(all_bg))
        n_pixels   = (y2 - y) * (x2 - x)
        background_signal = bg_median * n_pixels
    else:
        background_signal = 0.0

    adjusted_signal = max(0.0, raw_signal - background_signal)

    return {
        "rawSignal": round(raw_signal, 2),
        "backgroundSignal": round(background_signal, 2),
        "adjustedSignal": round(adjusted_signal, 2),
        "x": x, "y": y, "w": w, "h": h,
    }


def validate_box(box):
    if not isinstance(box, dict):
        raise PublicError("Each box must be an object.")
    for key in ("x", "y", "w", "h"):
        value = box.get(key)
        if not isinstance(value, (int, float)) or not np.isfinite(value):
            raise PublicError(f"Box {key} must be a finite number.")
    if box["w"] <= 0 or box["h"] <= 0:
        raise PublicError("Box width and height must be positive.")
    if box["w"] * box["h"] > MAX_IMAGE_PIXELS:
        raise PublicError("A selected box is too large.")


@app.route("/cleanup", methods=["POST"])
@app.route("/api/cleanup", methods=["POST"])
@limiter.limit("30 per minute")
def cleanup_temp_files():
    data = request.get_json(silent=True) or {}
    session_id = request_session_id(data)
    descriptors = []
    descriptors.extend(collect_blot_file_descriptors(data.get("blots", [])))
    if isinstance(data.get("files"), list):
        descriptors.extend(data["files"])
    if isinstance(data.get("upload"), dict):
        descriptors.append(data["upload"])
    delete_temp_files(descriptors, session_id, allow_uploads=True)
    return jsonify({"status": "deleted"})


# ─── PowerPoint generation ────────────────────────────────────────────────────

@app.route("/generate-pptx", methods=["POST"])
@app.route("/api/generate-pptx", methods=["POST"])
@limiter.limit("10 per minute")
def generate_pptx():
    try:
        data        = request.get_json(silent=True) or {}
        slides_data = data.get("slides", [])
        validate_pptx_payload(slides_data)

        prs              = Presentation()
        prs.slide_width  = Inches(10)
        prs.slide_height = Inches(7.5)
        blank_layout     = prs.slide_layouts[6]

        # Separate image slides from graph slides
        image_slide_data = next((s for s in slides_data if s.get("type") == "images"), None)
        graph_slides     = [s for s in slides_data if s.get("type") == "graphs"]

        # Build image slide first
        if image_slide_data:
            slide = prs.slides.add_slide(blank_layout)
            build_image_slide(slide, image_slide_data, prs)

        # Build graph slides, 4 per slide
        for slide_data in graph_slides:
            graphs   = slide_data.get("graphs", [])
            title    = slide_data.get("title", "")
            per_page = 4

            for page_start in range(0, max(1, len(graphs)), per_page):
                page_graphs = graphs[page_start:page_start + per_page]
                slide       = prs.slides.add_slide(blank_layout)
                build_graph_slide(slide, {
                    "title":  title if page_start == 0 else f"{title} (cont.)",
                    "graphs": page_graphs,
                }, prs)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(
            buf,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name="western-blot-analysis.pptx"
        )

    except Exception as error:
        app.logger.exception("PowerPoint generation failed")
        return error_response(error, "PowerPoint generation failed.")


def validate_pptx_payload(slides_data):
    if not isinstance(slides_data, list):
        raise PublicError("Slides must be a list.")
    if len(slides_data) > MAX_PPTX_SLIDES:
        raise PublicError(f"Too many slides. Maximum is {MAX_PPTX_SLIDES}.")

    graph_count = 0
    image_count = 0
    total_image_bytes = 0
    for slide in slides_data:
        if not isinstance(slide, dict):
            raise PublicError("Each slide must be an object.")
        images = slide.get("images", [])
        graphs = slide.get("graphs", [])
        if not isinstance(images, list) or not isinstance(graphs, list):
            raise PublicError("PowerPoint images and graphs must be lists.")
        image_count += len(images)
        if image_count > MAX_PPTX_SLIDES:
            raise PublicError("Too many blot images in PowerPoint export.")
        graph_count += len(graphs)
        for item in images:
            if not isinstance(item, dict):
                raise PublicError("Invalid blot image entry.")
            total_image_bytes += validate_data_url_size(item.get("image", ""))
            if len(str(item.get("label", ""))) > MAX_NAME_LENGTH:
                raise PublicError("A PowerPoint image label is too long.")
        for graph in graphs:
            total_image_bytes += validate_data_url_size(graph)
    if graph_count > MAX_PPTX_GRAPHS:
        raise PublicError(f"Too many graphs. Maximum is {MAX_PPTX_GRAPHS}.")
    if total_image_bytes > MAX_PPTX_TOTAL_IMAGE_BYTES:
        raise PublicError("PowerPoint image data exceeds the total size limit.", 413)


def validate_data_url_size(value):
    if not isinstance(value, str):
        raise PublicError("Invalid image data.")
    if not re.match(r"^data:image/(jpeg|jpg|png);base64,", value, re.IGNORECASE):
        raise PublicError("Export images must be JPEG or PNG base64 data URLs.")
    encoded = value.split(",", 1)[-1]
    approx_bytes = (len(encoded) * 3) // 4
    if approx_bytes > MAX_PPTX_IMAGE_BYTES:
        raise PublicError("An exported image is too large for PowerPoint.")
    return approx_bytes


def add_image_from_base64(slide, b64_string, left, top, width, height):
    img_bytes  = base64.b64decode(b64_string.split(",")[-1], validate=True)
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, left, top, width, height)


def add_label(slide, text, left, top, width, height, font_size=14, bold=False, align=PP_ALIGN.CENTER):
    text_box        = slide.shapes.add_textbox(left, top, width, height)
    tf              = text_box.text_frame
    tf.word_wrap   = True
    p              = tf.paragraphs[0]
    p.alignment    = align
    run            = p.add_run()
    run.text       = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = RGBColor(0x1f, 0x29, 0x33)


def build_image_slide(slide, slide_data, prs):
    images  = slide_data.get("images", [])
    if not images:
        return

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin  = Inches(0.4)
    label_h = Inches(0.3)
    spacing = Inches(0.2)

    add_label(slide, "Western Blot Images",
              margin, Inches(0.1), slide_w - margin * 2, Inches(0.35),
              font_size=18, bold=True)

    n               = len(images)
    total_label     = label_h * n
    total_spacing   = spacing * (n - 1)
    available_h     = slide_h - Inches(0.6) - total_label - total_spacing - margin
    full_img_w      = slide_w - margin * 2
    img_w           = full_img_w * 2 / 3
    img_h           = available_h / n
    left_x          = margin + (full_img_w - img_w) / 2

    current_y = Inches(0.6)
    for item in images:
        b64   = item.get("image")
        label = item.get("label", "")
        if b64:
            try:
                add_image_from_base64(slide, b64, left_x, current_y, img_w, img_h)
            except Exception:
                pass
        current_y += img_h
        add_label(slide, label, left_x, current_y, img_w, label_h,
                  font_size=12, bold=True, align=PP_ALIGN.CENTER)
        current_y += label_h + spacing


def build_graph_slide(slide, slide_data, prs):
    graphs  = slide_data.get("graphs", [])
    title   = slide_data.get("title", "")
    if not graphs:
        return

    slide_w  = prs.slide_width
    slide_h  = prs.slide_height
    margin   = Inches(0.4)
    padding  = Inches(0.2)

    add_label(slide, title,
              margin, Inches(0.1), slide_w - margin * 2, Inches(0.4),
              font_size=16, bold=True)

    cols    = 2
    rows    = 2
    img_w   = (slide_w - margin * 2 - padding * (cols - 1)) / cols
    img_h   = (slide_h - Inches(0.55) - margin - padding * (rows - 1)) / rows

    for i, b64 in enumerate(graphs):
        if i >= cols * rows:
            break
        col = i % cols
        row = i // cols
        x   = margin + col * (img_w + padding)
        y   = Inches(0.55) + row * (img_h + padding)
        if b64:
            try:
                add_image_from_base64(slide, b64, x, y, img_w, img_h)
            except Exception:
                pass

# ─── Run ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    debug = os.environ.get("FLASK_DEBUG", "false").lower() == "true"
    port  = int(os.environ.get("PORT", 5000))
    app.run(debug=debug, port=port)
